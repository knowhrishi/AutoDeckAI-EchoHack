# =========================================
# STEP 1: Imports & Basic Setup
# =========================================
import os
import re
import requests
import streamlit as st

from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_community.document_loaders import PyPDFLoader, PyMuPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.messages import HumanMessage
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import FAISS

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from unstructured.partition.pdf import partition_pdf
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
import base64

from chromadb.config import Settings
from langchain_core.documents import Document

from PyPDF2 import PdfReader
from faiss_vector_store import create_vectorstore

def extract_and_caption_pdf_elements(
    pdf_file_path: str,
    openai_api_key: str,
    output_dir: str = "./content/"
) -> list:
    """
    1) Extract figure images and possible table blocks using PyMuPDF (fitz).
    2) Generate short LLM-based captions for each figure/table.
    3) Return a list of dictionaries, each containing:
       {
         "type": "figure" or "table",
         "file_path": ...,
         "static_path": ...,
         "caption": ...,
         "figure_number": ...
       }
    """
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs("static", exist_ok=True)
    
    all_results = []
    
    # Install PyMuPDF if missing
    try:
        import pymupdf
        import pymupdf as fitz
    except ImportError:
        print(ImportError)
        # print("Installing PyMuPDF...")
        # pip.main(['install', '--quiet', '--no-cache-dir', 'PyMuPDF==1.21.1'])
        # import fitz
    print(fitz.__version__)

    try:
        print(f"PyMuPDF version: {fitz.version}")
        # Try opening the PDF
        try:
            doc = fitz.Document(pdf_file_path)
            print("Opened PDF with fitz.Document")
        except Exception as e1:
            try:
                doc = fitz.open(pdf_file_path)
                print("Opened PDF with fitz.open")
            except Exception as e2:
                raise Exception(f"Could not open PDF: {str(e1)}, {str(e2)}")
        
        # We will use a single LLM instance for captions
        llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.3,
            openai_api_key=openai_api_key
        )
        
        figure_counter = 1
        
        # Process each page for images + possible table blocks
        for page_num in range(len(doc)):
            page = doc[page_num]
            print(f"Processing page {page_num + 1} of {len(doc)}")
            
            # 1) Extract images
            image_list = page.get_images(full=True)
            print(f"Found {len(image_list)} images on page {page_num + 1}")
            
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        print(f"No image data for image {img_index} on page {page_num + 1}")
                        continue
                    
                    ext = base_image["ext"]
                    image_bytes = base_image["image"]
                    
                    # Save image
                    image_filename = f"figure_{page_num}_{img_index}.{ext}"
                    image_path = os.path.join(output_dir, image_filename)
                    static_path = os.path.join("static", image_filename)
                    
                    with open(image_path, "wb") as f_img:
                        f_img.write(image_bytes)
                    with open(static_path, "wb") as f_img:
                        f_img.write(image_bytes)
                    
                    print(f"Saved image to {image_path} and {static_path}")
                    
                    # Surrounding text for context
                    surrounding_text = page.get_text()
                    
                    # Generate a short caption with LLM
                    caption_chain = (
                        ChatPromptTemplate.from_template(
                            "Provide a short 1-2 sentence description of this image context:\n{text}"
                        )
                        | llm
                        | StrOutputParser()
                    )
                    caption = caption_chain.invoke({"text": surrounding_text})
                    
                    result = {
                        "type": "figure",
                        "file_path": image_path,
                        "static_path": static_path,
                        "caption": caption,
                        "figure_number": figure_counter
                    }
                    figure_counter += 1
                    all_results.append(result)
                    print(f"Added figure to results (figure_number={result['figure_number']})")
                except Exception as e:
                    print(f"Error processing image {img_index} on page {page_num+1}: {str(e)}")
            
            # 2) Basic "table" detection by text blocks
            text = page.get_text()
            page_text_blocks = text.split('\n\n')
            
            for block_index, block in enumerate(page_text_blocks):
                if is_likely_table(block):
                    try:
                        table_filename = f"table_{page_num}_{block_index}.txt"
                        table_path = os.path.join(output_dir, table_filename)
                        static_table_path = os.path.join("static", table_filename)
                        
                        # Save table content
                        with open(table_path, "w", encoding="utf-8") as f_tbl:
                            f_tbl.write(block)
                        with open(static_table_path, "w", encoding="utf-8") as f_tbl2:
                            f_tbl2.write(block)
                        
                        # Summarize the table
                        table_prompt = ChatPromptTemplate.from_template("""
                        Analyze this potential table content and provide a brief summary:
                        {content}
                        
                        If this is truly tabular data, summarize the key info.
                        If it's not really a table, respond with "Not a table."
                        """)
                        table_chain = table_prompt | llm | StrOutputParser()
                        caption = table_chain.invoke({"content": block})
                        
                        if "not a table" not in caption.lower():
                            # We'll treat it as a valid table
                            result = {
                                "type": "table",
                                "file_path": table_path,
                                "static_path": static_table_path,
                                "caption": caption,
                                "figure_number": figure_counter
                            }
                            figure_counter += 1
                            all_results.append(result)
                            print(f"Added table to results (figure_number={result['figure_number']})")
                    except Exception as e:
                        print(f"Error processing table block on page {page_num+1}: {str(e)}")
        
        doc.close()
        print(f"Extraction complete. Found {len(all_results)} items total.")
        return all_results
    except Exception as e:
        print(f"Error in PDF processing: {str(e)}")
        import traceback
        traceback.print_exc()
        return []

def is_likely_table(text: str) -> bool:
    """Heuristics to guess if a block is tabular data."""
    indicators = [
        text.count("|") > 2,
        text.count("\t") > 2,
        text.count("  ") > 8,
        bool(re.search(r"Table \d+[:.]\s", text, re.IGNORECASE)),
        len(re.findall(r"\d+\s+\d+\s+\d+", text)) > 0,
        bool(re.search(r"(\w+\s+){3,}\n(\w+\s+){3,}", text)),
        any(word in text.lower() for word in ["total", "sum", "average", "mean", "std", "min", "max"]),
        text.count("%") > 2,
        len(re.findall(r"\d+\.\d+", text)) > 3
    ]
    return sum(indicators) >= 2




# =========================================
# STEP 2: Streamlit UI & Inputs
# =========================================
st.set_page_config(page_title="Eco-centric Slide Generator", layout="centered")
st.title("🌿 Eco-centric Slide Generator")
st.markdown(
    """
    This tool converts **ecological research PDFs** (or via a DOI/URL) into **professionally formatted** PowerPoint slides.
    """
)
author_name = st.sidebar.text_input("Enter the author's name:")

openai_api_key = st.sidebar.text_input("Enter your OpenAI API key:", type="password")

presentation_focus = st.sidebar.selectbox(
    "Select the target audience or purpose of the presentation:",
    ["Researcher", "Practitioner", "Funding Body"]
)
num_slides = st.sidebar.number_input("Enter the number of slides to generate (including title slide):", min_value=1, value=7)

input_type = st.sidebar.radio("Select Input Type:", ["Upload PDF", "Enter DOI/URL"])
uploaded_file = None
doi_or_url = None

if input_type == "Upload PDF":
    uploaded_file = st.sidebar.file_uploader("📄 Upload a PDF document", type=["pdf"])
elif input_type == "Enter DOI/URL":
    doi_or_url = st.sidebar.text_input("🔗 Enter DOI or URL:")


# =========================================
# STEP 3: Utility Functions
# =========================================
def download_pdf_from_url(url: str) -> str:
    """Downloads a PDF from a given URL and returns the local file path."""
    response = requests.get(url)
    if response.status_code == 200:
        file_path = "downloaded_document.pdf"
        with open(file_path, "wb") as f:
            f.write(response.content)
        return file_path
    else:
        st.error("Failed to download PDF. Check the URL.")
        return ""

def extract_content_from_pdf(file_path: str) -> str:
    """
    Loads a PDF using PyPDFLoader (or PyMuPDFLoader) and
    concatenates all page_content into a single text string.
    """
    loader = PyPDFLoader(file_path)  # Or PyMuPDFLoader(file_path)
    documents = loader.load()
    combined_text = ""
    for doc in documents:
        combined_text += doc.page_content
    return combined_text

def preprocess_text_for_ecology(text: str) -> str:
    """
    Removes headers, footers, or references to
    clean the text for ecological summarization.
    """
    cleaned_text = re.sub(r"\nReferences.*", "", text, flags=re.IGNORECASE)
    cleaned_text = re.sub(r"\nPage \d+", "", cleaned_text)
    return cleaned_text

def parse_llm_response(response_content: str):
    """
    Parses the LLM response into a list of slide dicts with improved error handling.
    """
    if not response_content or not isinstance(response_content, str):
        print("Invalid response content")
        return generate_default_slide_list()
        
    slides = []
    try:
        # Clean up response content
        response_content = response_content.strip()
        
        # Split into lines and remove empty lines
        lines = [line.strip() for line in response_content.split('\n') if line.strip()]
        
        current_slide = None
        content_buffer = []
        
        for line in lines:
            # Skip separator lines and empty lines
            if line == '---' or not line:
                continue
                
            # Remove markdown formatting
            line = line.replace('**', '')
            
            if 'Title:' in line and line.lower().startswith('slide'):
                # Save previous slide if it exists
                if current_slide:
                    current_slide['content'] = '\n'.join(content_buffer).strip()
                    slides.append(current_slide)
                
                # Start new slide
                slide_title = line.split('Title:', 1)[1].strip()
                current_slide = {'title': slide_title}
                content_buffer = []
                
            elif 'Content:' in line and line.lower().startswith('slide'):
                content = line.split('Content:', 1)[1].strip()
                if content:
                    content_buffer.append(content)
            elif current_slide is not None:
                # Append any other lines to content buffer
                content_buffer.append(line)
        
        # Add the last slide
        if current_slide:
            current_slide['content'] = '\n'.join(content_buffer).strip()
            slides.append(current_slide)
        
        # Validate final result
        if not slides:
            print("No valid slides parsed")
            return generate_default_slide_list()
            
        # Debug print
        print(f"Successfully parsed {len(slides)} slides")
        # for slide in slides:
        #     print(f"Slide: {slide['title']}")
            
        return slides
        
    except Exception as e:
        print(f"Error parsing LLM response: {str(e)}")
        return generate_default_slide_list()
    
def generate_default_slide_list():
    """
    Generates a default list of slides
    """
    return [
        {"title": "Document Overview", "content": "- Key points extracted from the document"},
        {"title": "Main Findings", "content": "- Important findings and insights\n- Key takeaways from the text"},
        {"title": "References", "content": "- Document sources and citations"},
        {"title": "Thank You", "content": "Thank you for your attention!"}
    ]

def remove_slide(prs, slide_index):
    """
    Removes a slide from the presentation by its index.
    """
    slide_id = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[slide_index]


def update_slide_content_with_figures(content: str, element_lookup: dict) -> str:
    """Helper function to update slide content with correct figure references"""
    updated_content = content
    
    # Create a mapping of figure numbers to elements
    number_to_elem = {
        elem['figure_number']: elem 
        for elem in element_lookup.values()
    }
    
    # Replace figure references with full paths
    for num in range(1, len(number_to_elem) + 1):
        figure_pattern = f"[FIGURE {num}]"
        table_pattern = f"[TABLE {num}]"
        
        if figure_pattern in content or table_pattern in content:
            if num in number_to_elem:
                elem = number_to_elem[num]
                if elem['type'].lower() == 'figure':
                    updated_content = updated_content.replace(
                        figure_pattern,
                        f"{figure_pattern} {elem['file_path']}"
                    )
                else:
                    updated_content = updated_content.replace(
                        table_pattern,
                        f"{table_pattern} {elem['file_path']}"
                    )
    
    return updated_content


def generate_presentation(slides: list, author_name: str, extracted_elements: list) -> str:
    """Creates a PowerPoint presentation with integrated figures and tables."""
    try:
        # Create directories if they don't exist
        os.makedirs("static", exist_ok=True)
        
        # Load template
        prs = Presentation('autodeckai2.pptx')
        
        # Remove existing slides
        for i in range(len(prs.slides) - 1, -1, -1):
            remove_slide(prs, i)
        
        # Create element lookup with both file path and figure number
        element_lookup = {}
        for elem in extracted_elements:
            element_lookup[f"[FIGURE {elem['figure_number']}]"] = elem
            element_lookup[f"[TABLE {elem['figure_number']}]"] = elem
        
        # Layout indices
        TITLE_SLIDE_LAYOUT = 0
        CONTENT_SLIDE_LAYOUT = 1
        REFERENCES_SLIDE_LAYOUT = 2
        THANK_YOU_SLIDE_LAYOUT = 3
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = slides[0].get('title', 'Presentation Title')
        subtitle.text = f"Author: {author_name}"
        
        # Content Slides
        main_slides = slides[1:-2] if len(slides) > 3 else slides[1:2]
        
        for slide_data in main_slides:
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_SLIDE_LAYOUT])
            
            # Add title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Untitled Slide')
            
            # Process content
            content = slide.placeholders[1].text_frame
            content.clear()
            
            current_y = Inches(2)
            content_lines = slide_data.get('content', '').split('\n')
            
            # First pass: Add text content and track figure references
            text_content = []
            figures_to_add = []
            
            for line in content_lines:
                line = line.strip()
                if not line:
                    continue
                
                # Check for figure/table references
                for ref, elem in element_lookup.items():
                    if ref in line:
                        figures_to_add.append(elem)
                        # Remove the reference from the line
                        line = line.replace(ref, '')
                
                if line:  # Only add non-empty lines
                    text_content.append(line)
            
            # Add text content first
            for line in text_content:
                p = content.add_paragraph()
                if line.startswith('*') or line.startswith('-'):
                    p.bullet = True
                    p.text = line.lstrip('*- ').strip()
                else:
                    p.text = line
                p.font.size = Pt(18)
            
            # Then add figures below the text
            current_y = Inches(4)  # Start figures below text content
            for elem in figures_to_add:
                try:
                    if os.path.exists(elem['static_path']):
                        # Add figure
                        img_width = Inches(6)
                        img_left = Inches(1.5)
                        
                        picture = slide.shapes.add_picture(
                            elem['static_path'],
                            img_left,
                            current_y,
                            width=img_width
                        )
                        
                        # Calculate height based on aspect ratio
                        actual_height = picture.height
                        
                        # Add caption below image
                        caption_top = current_y + Inches(actual_height/914400)  # Convert EMUs to inches
                        caption_box = slide.shapes.add_textbox(
                            img_left,
                            caption_top,
                            img_width,
                            Inches(0.5)
                        )
                        caption_para = caption_box.text_frame.add_paragraph()
                        caption_para.text = f"{elem['type'].title()} {elem['figure_number']}: {elem['caption']}"
                        caption_para.font.size = Pt(12)
                        caption_para.font.italic = True
                        
                        current_y = caption_top + Inches(0.7)
                        print(f"Added {elem['type']} {elem['figure_number']}")
                    else:
                        print(f"File not found: {elem['static_path']}")
                        
                except Exception as e:
                    print(f"Error adding figure: {str(e)}")
                    continue
                # Add text line
            p = content.add_paragraph()
            if line.startswith('-'):
                p.bullet = True
                p.text = line.lstrip('-').strip()
            else:
                p.text = line
            p.font.size = Pt(18)
        
    
        # == C. References Slide ==
        ref_slide = prs.slides.add_slide(prs.slide_layouts[REFERENCES_SLIDE_LAYOUT])
        ref_title = ref_slide.shapes.title
        # Safely get references slide content
        if len(slides) >= 3:
            ref_title.text = slides[-2].get('title', 'References')
            references_text = slides[-2].get('content', 'No references available.')
        else:
            ref_title.text = 'References'
            references_text = 'No references available.'
            
        ref_content = ref_slide.placeholders[1]
        ref_frame = ref_content.text_frame
        ref_frame.clear()
        
        ref_lines = references_text.split('\n')
        
        # Add references
        for line in ref_lines:
            line = line.strip()
            if not line:
                continue
            p = ref_frame.add_paragraph()
            p.text = line
            p.bullet = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.alignment = PP_ALIGN.LEFT
        
        # Add figure/table attributions using lookup
        if element_lookup:
            p = ref_frame.add_paragraph()
            p.text = "Figure and Table Sources:"
            p.font.bold = True
            p.font.size = Pt(16)
            
            for path, elem in element_lookup.items():
                p = ref_frame.add_paragraph()
                p.text = f"{elem['type'].title()}: {elem['caption']}"
                p.bullet = True
                p.font.size = Pt(14)
        
        # == D. Thank You Slide ==
        thanks_slide = prs.slides.add_slide(prs.slide_layouts[THANK_YOU_SLIDE_LAYOUT])
        thanks_title = thanks_slide.shapes.title    
        
        # Safely get thank you slide content
        if len(slides) >= 4:
            thanks_title.text = slides[-1].get('title', 'Thank You')
            thanks_content = slides[-1].get('content', 'We appreciate your attention!')
        else:
            thanks_title.text = 'Thank You'
            thanks_content = 'We appreciate your attention!'
        
        if len(thanks_slide.placeholders) > 1:
            thanks_subtitle = thanks_slide.placeholders[1]
            thanks_subtitle.text = thanks_content    
        if len(thanks_slide.placeholders) > 1:
            thanks_subtitle = thanks_slide.placeholders[1]
            thanks_subtitle.text = slides[-1].get('content', 'We appreciate your attention!')
        
        # Save the presentation
        output_filename = "generated_presentation.pptx"
        prs.save(output_filename)
        print("Presentation saved successfully")
        return output_filename
        
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        return ""

st.sidebar.write("---")
# For demonstration:
st.sidebar.markdown("🛠️ Only for testing purpose")
if st.sidebar.button("Create Demo Slides"):
    # add markdown saying that it is for testing purposes
    sample_slides = [
        {"title": "Introduction", "content": ""},
        {"title": "Overview", "content": "- Purpose\n- Scope\n- Approach"},
        {"title": "Results", "content": "- Observed data\n- Statistical insights\n\n- Graphical analysis"},
        {"title": "References", "content": "- Smith et al. 2020\n- Doe and Roe, 2019"},
        {"title": "Thank You", "content": "Feel free to reach out with any questions!"}
    ]
    file_path = generate_presentation(sample_slides, author_name="Jane Doe")
    st.success("Presentation generated!")
    with open(file_path, "rb") as f:
        st.download_button("Download PPTX", f.read(), "generated_presentation.pptx")



def generate_slides_with_retrieval(vectorstore, presentation_focus: str, num_slides: int, extracted_elements: list, openai_api_key: str):
    """
    Uses a RetrievalQA chain to combine retrieved content into slides.
    Now with improved response handling and debugging.
    """
    try:
        retriever = vectorstore.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 10}
        )    

        figures_info = "\nAvailable Figures and Tables:\n"
        for elem in extracted_elements:
            figures_info += f"- {elem['type'].title()}: {elem['caption']}\n"
        
        prompt_text = (
            f"As a **{presentation_focus}**, create a slide presentation with **{num_slides}** total slides "
            "using the content provided. You have access to the following figures and tables that you can incorporate:\n"
            f"{figures_info}\n"
            "Your presentation must include:\n\n"
            "1. **Title Page** (Slide 1):\n"
            "   - Only the main title (paper or project name) and author name(s).\n\n"
            "2. **Main Slides** (Slides 2 through N-2):\n"
            "   - Each slide has:\n"
            "       - A clear, descriptive title (e.g., 'Methodology', 'Results', etc.)\n"
            "       - Bullet-pointed content that summarizes key points\n"
            "       - Where appropriate, include [FIGURE X] or [TABLE X] markers\n"
            "   - When referencing figures or tables, integrate them naturally into the content\n"
            "   - Ensure proper attribution for any included figures\n\n"
            "3. **Conclusion** (Slide N-1):\n"
            "   - Summarize main findings or takeaways\n"
            "   - Include any recommendations\n\n"
            "4. **References** (Slide N):\n"
            "   - List sources and include attributions for any used figures\n\n"
            "5. **Thank You** (Final Slide):\n"
            "   - Brief closing message\n\n"
            "Format each slide EXACTLY as follows:\n"
            "Slide 1 Title: [Title]\n"
            "Slide 1 Content: [Content]\n"
            "Slide 2 Title: [Title]\n"
            "Slide 2 Content: [Content]\n"
            "... and so on for each slide.\n\n"
            "Important: Keep this exact format for each slide.\n"
        )

        chain = RetrievalQA.from_chain_type(
            llm=ChatOpenAI(
                openai_api_key=openai_api_key,
                model_name="gpt-4o",
                temperature=0.7,
                max_tokens=2000
            ),
            retriever=retriever,
            chain_type="stuff",
            return_source_documents=True
        )
        
        # Execute the chain
        response = chain.invoke({"query": prompt_text})
        
        # Handle different response formats
        if isinstance(response, dict):
            if 'result' in response:
                result = response['result']
            elif 'answer' in response:
                result = response['answer']
            else:
                result = str(response)
        else:
            result = str(response)

        # # Debug print
        # print("Raw LLM Response:", result)
        
        # Clean up the response
        result = result.strip()
        if result.startswith('Here\'s') or result.startswith('---'):
            # Remove any prefix text before the first slide
            first_slide_idx = result.find('Slide 1 Title:')
            if first_slide_idx != -1:
                result = result[first_slide_idx:]
        
        # Validate the response format more carefully
        if "Slide 1 Title:" not in result:
            print("Response missing required 'Slide 1 Title:' format")
            return generate_default_slides()

        return result
        
    except Exception as e:
        print(f"Error in slide generation: {str(e)}")
        return generate_default_slides()

def generate_default_slides():
    """
    Generates a default set of slides when the main generation fails
    """
    return """
    Slide 1 Title: Document Overview
    Slide 1 Content: - Key points extracted from the document
    
    Slide 2 Title: Main Findings
    Slide 2 Content: - Important findings and insights
    - Key takeaways from the text
    
    Slide 3 Title: References
    Slide 3 Content: - Document sources and citations
    
    Slide 4 Title: Thank You
    Slide 4 Content: Thank you for your attention!
    """
# =========================================
# STEP 4: Main Logic
# =========================================

st.write("---")
generate_slides_clicked = st.button("Generate Slide Deck")
if generate_slides_clicked:
    if not openai_api_key:
        st.error("Please provide a valid OpenAI API key.")
    elif not (uploaded_file or doi_or_url):
        st.error("Please upload a PDF or provide a DOI/URL.")
    else:

        status_placeholder = st.empty()

        progress_bar = st.progress(0)
        status_placeholder.info("Processing your input...")

        # A. Download or store PDF
        file_path = ""
        if uploaded_file:
            file_path = "uploaded_document.pdf"
            with open(file_path, "wb") as f:
                f.write(uploaded_file.read())
        elif doi_or_url:
            file_path = download_pdf_from_url(doi_or_url)

        if file_path:
            progress_bar.progress(25)
            status_placeholder.info("Extracting and cleaning text...")

            extracted_text = extract_content_from_pdf(file_path)
            cleaned_text = preprocess_text_for_ecology(extracted_text)

            progress_bar.progress(50)
            status_placeholder.info("Creating/Loading vector store...")


            vectorstore = create_vectorstore(cleaned_text, openai_api_key)

                # 3) NEW: Extract & caption PDF images/tables
            extracted_elements = extract_and_caption_pdf_elements(
                pdf_file_path=file_path,
                openai_api_key=openai_api_key,
                output_dir="content/"
            )

            progress_bar.progress(70)
            status_placeholder.info("Generating slides via LLM retrieval...")

            llm_response = generate_slides_with_retrieval(vectorstore, presentation_focus, num_slides, extracted_elements, openai_api_key)
            slides = parse_llm_response(llm_response)

            progress_bar.progress(90)
            status_placeholder.info("Creating PowerPoint presentation...")

            pptx_file = generate_presentation(slides, author_name, extracted_elements)

            progress_bar.progress(100)
            status_placeholder.success("🎉 Slides generated successfully!")
            st.download_button(
                label="📥 Download Presentation",
                data=open(pptx_file, "rb").read(),
                file_name="EcoHack_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

            # Preview
            st.markdown("### 📄 Generated Slides Preview:")
            for slide in slides:
                st.markdown(f"**{slide['title']}**")
                st.write(slide['content'])
        else:
            st.warning("Unable to process the file. Please verify your input.")
else:
    st.info("Configure your inputs, then click 'Generate Slide Deck' to proceed.")

