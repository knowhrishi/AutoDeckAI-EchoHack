# utils.py
from pathlib import Path
from openai import OpenAI
# from markitdown import MarkItDown
import os, time
import re
import requests
import traceback
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_openai import ChatOpenAI
from langchain.chains import RetrievalQA
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_community.document_loaders import PyPDFLoader
from functools import lru_cache

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from typing import Set
import nltk
from nltk.corpus import stopwords
from sklearn.feature_extraction.text import TfidfVectorizer
from transformers import pipeline, BlipProcessor, BlipForConditionalGeneration
from pydantic import BaseModel, Field
from typing import List, Optional

# Eco-themed loading messages
ECO_LOADING_MESSAGES = [
    "🌱 Growing your ideas into slides...",
    "🦋 Pollinating content across sections...",
    "🌳 Branching out through your research...",
    "🍃 Cultivating sustainable visuals...",
    "🌿 Nurturing your presentation...",
    "🌸 Blossoming your insights...",
    "🐝 Cross-pollinating concepts...",
    "🌍 Eco-system analysis in progress..."
]
# Ecological term validation
ECOLOGICAL_TERMS = {
    'biodiversity', 'ecosystem', 'sustainability', 'carbon sequestration',
    'habitat', 'conservation', 'climate change', 'species richness',
    'ecological footprint', 'restoration', 'biome', 'keystone species'
}
# Eco tips
ECO_TIPS = [
    "📚 Did you know? Digital presentations save approximately 2-3 trees per year compared to paper handouts!",
    "💡 Using dark mode in presentations can reduce energy consumption by up to 60% on OLED screens.",
    "🌱 Your digital presentation contributes to reducing paper waste and promoting sustainable practices."
]
def get_image_captioner(model_name, api_key=None, model_provider="Hugging Face"):
    """Get image captioning based on provider"""
    if model_provider == "OpenAI":
        return {"type": "openai", "api_key": api_key}
    else:
        from transformers import pipeline
        return pipeline("image-to-text", model="Salesforce/blip-image-captioning-base", device="cpu")



def caption_image(raw_img, captioner, model_provider="Hugging Face"):
    """Generate caption using appropriate model"""
    try:
        if model_provider == "OpenAI":
            import base64
            from io import BytesIO
            from openai import OpenAI

            # Convert PIL image to base64
            buffered = BytesIO()
            raw_img.save(buffered, format="JPEG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()

            client = OpenAI(api_key=captioner["api_key"])
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{
                    "role": "user", 
                    "content": [
                        {"type": "text", "text": "Describe this scientific figure in one line."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_base64}"}}
                    ]
                }],
                max_tokens=100
            )
            return response.choices[0].message.content
        else:
            outputs = captioner(raw_img)
            return outputs[0]["generated_text"]
    except Exception as e:
        print(f"Caption generation failed: {str(e)}")
        return "Figure from document"
    
def download_pdf_from_url(url: str) -> str:
    """Downloads a PDF from a given URL and returns the local file path."""
    response = requests.get(url)
    if response.status_code == 200:
        file_path = "downloaded_document.pdf"
        with open(file_path, "wb") as f:
            f.write(response.content)
        return file_path
    return ""

def extract_content_from_pdf(file_path: str) -> str:
    """Loads a PDF and concatenates all page_content into a single text string."""
    loader = PyPDFLoader(file_path)
    documents = loader.load()
    return "".join(doc.page_content for doc in documents)

# Enhanced file processing
def extract_content_from_file(file_path: str) -> str:
    """Handle multiple file types"""
    from pptx import Presentation
    from docx import Document
    
    try:
        if file_path.endswith('.pdf'):
            return extract_content_from_pdf(file_path)
        
        elif file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        elif file_path.endswith('.docx'):
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif file_path.endswith('.pptx'):
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        return ""
    except Exception as e:
        print(f"Error processing {file_path}: {str(e)}")
        return ""

def preprocess_text_for_ecology(text: str) -> str:
    """Removes headers, footers, or references to clean the text."""
    # cleaned_text = re.sub(r"\nReferences.*", "", text, flags=re.IGNORECASE)
    cleaned_text = re.sub(r"\nPage \d+", "", text)
    return cleaned_text

def parse_llm_response(response_content: str) -> List[Dict[str, str]]:
    """Parses the LLM response into a list of slide dictionaries."""
    if not response_content or not isinstance(response_content, str):
        print("Invalid response content")
        return generate_default_slide_list()
        
    slides = []
    try:
        lines = [line.strip() for line in response_content.split('\n') if line.strip()]
        current_slide = None
        content_buffer = []
        
        for line in lines:
            if line == '---' or not line:
                continue
                
            line = line.replace('**', '')
            
            if 'Title:' in line and line.lower().startswith('slide'):
                if current_slide:
                    current_slide['content'] = '\n'.join(content_buffer).strip()
                    slides.append(current_slide)
                
                slide_title = line.split('Title:', 1)[1].strip()
                current_slide = {'title': slide_title}
                content_buffer = []
                
            elif 'Content:' in line and line.lower().startswith('slide'):
                content = line.split('Content:', 1)[1].strip()
                if content:
                    content_buffer.append(content)
            elif current_slide is not None:
                content_buffer.append(line)
        
        if current_slide:
            current_slide['content'] = '\n'.join(content_buffer).strip()
            slides.append(current_slide)
            
        return slides if slides else generate_default_slide_list()
        
    except Exception as e:
        print(f"Error parsing LLM response: {str(e)}")
        return generate_default_slide_list()

def generate_default_slide_list() -> List[Dict[str, str]]:
    """Generates a default list of slides."""
    return [
        {"title": "Document Overview", "content": "- Key points extracted from the document"},
        {"title": "Main Findings", "content": "- Important findings and insights\n- Key takeaways from the text"},
        {"title": "References", "content": "- Document sources and citations"},
        {"title": "Thank You", "content": "Thank you for your attention!"}
    ]

def clean_static_directory():
    """Clean the static directory before processing new files."""
    import shutil
    try:
        for dir_name in ["static", "content"]:
            if os.path.exists(dir_name):
                shutil.rmtree(dir_name)
            os.makedirs(dir_name)
        print("Successfully cleaned static and content directories")
    except Exception as e:
        print(f"Error cleaning directories: {str(e)}")

def generate_slides_with_retrieval(
    vectorstore: Any,
    presentation_focus: str,
    num_slides: int,
    extracted_elements: List[Dict[str, Any]],
    model_provider: str,  # Add this parameter
    model_name: str,       # Add this parameter
    api_key: str,           # Add this parameter,
    # ecological_theme: str,
    data_visualization: List[str]

) -> str:
    """Generates slides using retrieval-based approach."""
    try:
        retriever = vectorstore.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 10}
        )

        # Separate figures and tables for clarity
        figures_list = [elem for elem in extracted_elements if elem['type'].lower() == 'figure']
        tables_list = [elem for elem in extracted_elements if elem['type'].lower() == 'table']

        # Create a more explicit prompt for figures and tables
        figures_info = "\\nAvailable Figures:\\n"
        for fig in figures_list:
            marker = f"[FIGURE {fig['figure_number']}]"
            figures_info += f"- {marker}: {fig['caption']}\\n"

        tables_info = "\\nAvailable Tables:\\n"
        for tab in tables_list:
            marker = f"[TABLE {tab['figure_number']}]"
            tables_info += f"- {marker}: {tab['caption']}\\n"

        prompt_text = create_slide_prompt(
            presentation_focus, num_slides, 
            figures_info, 
            tables_info, 
            # ecological_theme, 
            data_visualization)
        llm = get_llm(model_provider, model_name, api_key)

        chain = RetrievalQA.from_chain_type(
            # llm=ChatOpenAI(
            #     openai_api_key=openai_api_key,
            #     model_name="gpt-4o",
            #     temperature=0.7,
            #     max_tokens=2000
            # ),
            llm=llm,
            retriever=retriever,
            chain_type="stuff",
            return_source_documents=True
        )
        
        response = chain.invoke({"query": prompt_text})
        result = response.get('result', response.get('answer', str(response))) if isinstance(response, dict) else str(response)
        
        result = clean_llm_response(result)
        
        return result if "Slide 1 Title:" in result else generate_default_slides()
        
    except Exception as e:
        print(f"Error in slide generation: {str(e)}")
        return generate_default_slides()

def create_slide_prompt(
        presentation_focus: str, 
        num_slides: int, 
        figures_info: str, 
        table_info: str,
        # ecological_theme: str, 
        data_visualization: list) -> str:
    """Creates an improved prompt for structured slide generation."""
    return (
        f"As an ecology/environmental science **{presentation_focus}**, create a **{num_slides}-slide** presentation focusing on sustainable and eco-friendly research insights. "
        "Each visual element should highlight ecological impacts, biodiversity considerations, and environmental sustainability.\n\n"
        "AVAILABLE RESEARCH VISUALS:\n"
        "FIGURES AVAILABLE:\n"
        f"{figures_info}\n\n"
        "TABLES AVAILABLE:\n"
        f"{table_info}\n\n"
        " 1. **Theme Integration**: \n"
        # f"Apply {ecological_theme} visual metaphors\n"
        f" Use {', '.join(data_visualization)} visualizations where appropriate\n"
        "**Structure Requirements:**\n\n"
        "1. **Title Slide** (Slide 1):\n"
        "   - Clean layout with ONLY:\n"
        "   - Bold presentation title (8 words max)\n"
        "   - Author/team name(s) centered\n\n"
        "2. **Content Slides** (Slides 2-{num_slides-3}):\n"
        "   - **Title:** Action-oriented header (e.g., 'Market Trends: 2023 Analysis')\n"
        "   - **Body:** 3-5 concise bullet points per slide:\n"
        "     - Lead with insights, not data\n"
        "     - Reference visuals using EXACT tags on separate lines\n"
        "     - Example:\n"
        "       *'45% revenue growth in Q3 (see regional breakdown)*\n"
        "       [FIGURE 2]\n\n"
        "3. **Conclusion Slide** (Slide {num_slides-2}):\n"
        "   - 3 key takeaways using data-driven language\n"
        "   - 1-2 actionable recommendations\n"
        "   - Optional closing visual: [FIGURE X] or [TABLE Y]\n\n"
        "4. **References** (Slide {num_slides-1}):\n"
        "   - APA/MLA format sources\n"
        "   - Complete figure/table credits\n\n"
        "5. **Closing Slide** (Slide {num_slides}):\n"
        "   - Minimal text: 'Key Questions?' or 'Next Steps'\n"
        "   - Optional: Contact info/logo\n\n"
        "**Critical Rules:**\n"
        "- Place EVERY [FIGURE X]/[TABLE Y] reference on its own line after explanatory text\n"
        "- Maintain narrative flow: Problem → Analysis → Solution\n"
        "- Use {presentation_focus} terminology appropriately\n"
        "- Never exceed {num_slides} slides\n\n"
        "Format EXACTLY like:\n"
        "Slide 1 Title: [Innovation Report 2024]\n"
        "Slide 1 Content: [Dr. Jane Smith]\n"
        "Slide 2 Title: [Emerging Market Patterns]\n"
        "Slide 2 Content:\n"
        "- Consumer tech adoption up 62% since 2020\n"
        "[FIGURE 1]\n"
        "- Regional variance exceeds predictions\n"
        "[TABLE 3]"
    )

def clean_llm_response(result: str) -> str:
    """Cleans up the LLM response."""
    result = result.strip()
    if result.startswith('Here\'s') or result.startswith('---'):
        first_slide_idx = result.find('Slide 1 Title:')
        if first_slide_idx != -1:
            result = result[first_slide_idx:]
    return result

def generate_default_slides() -> str:
    """Generates default slides when the main generation fails."""
    return """
    Slide 1 Title: Document Overview
    Slide 1 Content: - Key points extracted from the document
    
    Slide 2 Title: Main Findings
    Slide 2 Content: - Important findings and insights
    - Key takeaways from the text
    
    Slide 3 Title: References
    Slide 3 Content: - Document sources and citations
    
    Slide 4 Title: Thank You
    Slide 4 Content: Thank you for your attention!
    """

# PDF Element Processing Functions
def is_likely_table(text: str) -> bool:
    """
    Rudimentary check if the text block might be tabular data
    using certain heuristics like repeated columns or numeric patterns.
    """
    indicators = [
        text.count("|") > 2,
        text.count("\t") > 2,
        text.count("  ") > 8,
        bool(re.search(r"Table \d+[:.]\s", text, re.IGNORECASE)),
        len(re.findall(r"\d+\s+\d+\s+\d+", text)) > 0,
        bool(re.search(r"(\w+\s+){3,}\n(\w+\s+){3,}", text)),
        any(word in text.lower() for word in ["total", "sum", "average", "mean", "std", "min", "max"]),
        text.count("%") > 2,
        len(re.findall(r"\d+\.\d+", text)) > 3
    ]
    return sum(indicators) >= 2


def is_valid_figure(base_image: dict, min_size: int = 10000) -> bool:
    """Determines if an image is likely to be a meaningful figure."""
    try:
        image_bytes = base_image["image"]
        if len(image_bytes) < min_size:
            return False
        
        width = base_image.get("width", 0)
        height = base_image.get("height", 0)
        
        if width > 0 and height > 0:
            # Filter out extremely small or extremely large
            if width < 100 or height < 100 or width > 3000 or height > 3000:
                return False
        
        ext = base_image["ext"].lower()
        return ext in {"jpg", "jpeg", "png", "bmp"}
    except Exception as e:
        print(f"Error checking image validity: {str(e)}")
        return False

def get_surrounding_text(page: Any, rect: Any, buffer: int = 50) -> str:
    """Get text surrounding an image with a buffer zone."""
    import pymupdf as fitz
    try:
        x0 = max(0, rect.x0 - buffer)
        y0 = max(0, rect.y0 - buffer)
        x1 = min(page.rect.width, rect.x1 + buffer)
        y1 = min(page.rect.height, rect.y1 + buffer)
        
        clip_rect = fitz.Rect(x0, y0, x1, y1)
        return page.get_text("text", clip=clip_rect)
    except Exception as e:
        print(f"Error getting surrounding text: {str(e)}")
        return page.get_text()

# Presentation Generation Functions
def remove_slide(prs: Presentation, slide_index: int):
    """Removes a slide from the presentation by its index."""
    slide_id = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[slide_index]

def update_slide_content_with_figures(content: str, element_lookup: Dict[str, Any]) -> str:
    """
    If you need to replace placeholders [FIGURE 1] with actual file paths, do it here.
    (Currently not used in this code, but can be helpful in certain flows.)
    """
    updated_content = content
    number_to_elem = {elem['figure_number']: elem for elem in element_lookup.values()}
    for num, elem in number_to_elem.items():
        figure_pattern = f"[FIGURE {num}]"
        table_pattern = f"[TABLE {num}]"
        if figure_pattern in content or table_pattern in content:
            pattern = figure_pattern if elem['type'].lower() == 'figure' else table_pattern
            updated_content = updated_content.replace(pattern, f"{pattern} {elem['file_path']}")
    return updated_content

async def generate_presentation(
    slides: list, 
    author_name: str, 
    extracted_elements: list, 
    text_llm,
    theme_file: str
) -> str:
    """
    Build the final PPTX, ensuring images and tables fit.
    If they don't fit, we either scale or move to a new slide.
    """
    from pptx import Presentation
    import traceback

    try:
        # Build a lookup for figure/table references
        element_lookup = {}
        for elem in extracted_elements:
            key = f"[{elem['type'].upper()} {elem['figure_number']}]"
            element_lookup[key] = elem
            print(f"Added {key} to lookup with path: {elem['static_path']}")

        prs = Presentation(theme_file)

        # Remove existing slides from template
        for i in range(len(prs.slides) - 1, -1, -1):
            remove_slide(prs, i)

        # Adjust these to match your template
        TITLE_SLIDE_LAYOUT = 0
        CONTENT_SLIDE_LAYOUT = 1
        REFERENCES_SLIDE_LAYOUT = 2
        THANK_YOU_SLIDE_LAYOUT = 3

        # 1) Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]  # May need try/except if no subtitle placeholder
        title.text = slides[0].get('title', 'Presentation Title')
        subtitle.text = f"Author: {author_name}"

        # We'll track a "current_y_in" for each new content slide
        # so we can stack images/tables below the bullet text
        main_slides = slides[1:-2] if len(slides) > 3 else slides[1:2]

        for slide_data in main_slides:
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_SLIDE_LAYOUT])

            # Slide Title
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get('title', 'Untitled Slide')

            # Slide body text
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            text_content = []
            elements_to_add = []

            # Parse lines
            lines = slide_data.get('content', '').split('\n')
            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # check for references
                found_ref = None
                for ref, elem in element_lookup.items():
                    if ref in line:
                        elements_to_add.append(elem)
                        # remove the reference text to avoid printing "[FIGURE X]" in bullet
                        line = line.replace(ref, '').strip()
                        found_ref = ref
                if line:
                    text_content.append(line)

            # Add bullet points
            for line in text_content:
                p = text_frame.add_paragraph()
                # simple bullet detection
                if line.startswith('-') or line.startswith('*'):
                    p.bullet = True
                    p.text = line.lstrip('-* ').strip()
                else:
                    p.text = line
                p.font.size = Pt(18)

            # We'll place images/tables below the text placeholder
            current_y_in = (
                (content_placeholder.top + content_placeholder.height) / 914_400.0 
                + 0.3  # small gap
            )

            # Insert each figure/table
            for elem in elements_to_add:
                try:
                    file_path = os.path.abspath(elem['static_path'])
                    if not os.path.exists(file_path):
                        print(f"File not found for {elem['type']} {elem['figure_number']}: {file_path}")
                        continue

                    if elem['type'].lower() == 'figure':
                        # Add image within bounds
                        print(f"Adding FIGURE {elem['figure_number']} from {file_path}")
                        slide, current_y_in = add_image_within_bounds(
                            prs,
                            slide,
                            file_path,
                            current_y_in,
                            left_in=1.0,
                            desired_width_in=6.0,
                            margin_bottom_in=0.3
                        )
                        # Then add a small caption text box
                        caption_textbox = slide.shapes.add_textbox(
                            Inches(1.0),
                            Inches(current_y_in),
                            Inches(6.0),
                            Inches(0.4)
                        )
                        caption_para = caption_textbox.text_frame.add_paragraph()
                        caption_para.text = f"Figure {elem['figure_number']}: {elem['caption']}"
                        caption_para.font.size = Pt(12)
                        caption_para.font.italic = True
                        current_y_in += 0.5

                    else:
                        # table
                        print(f"Adding TABLE {elem['figure_number']}")
                        # We'll pass 'slide' and see if we can place it
                        # The function returns new Y position
                        new_y = await add_formatted_table_element(slide, elem, current_y_in, text_llm)
                        current_y_in = new_y

                except Exception as e:
                    print(f"Error adding {elem['type']} {elem['figure_number']}: {str(e)}")
                    traceback.print_exc()
                    continue

        # 2) References slide
        ref_slide = prs.slides.add_slide(prs.slide_layouts[REFERENCES_SLIDE_LAYOUT])
        ref_slide.shapes.title.text = slides[-2].get('title', 'References')

        ref_text = slides[-2].get('content', '')
        ref_placeholder = ref_slide.placeholders[1]
        ref_frame = ref_placeholder.text_frame
        ref_frame.clear()

        for line in ref_text.split('\n'):
            line = line.strip()
            if line:
                p = ref_frame.add_paragraph()
                p.text = line
                p.bullet = True
                p.font.size = Pt(16)

        # 3) Thank You slide
        thanks_slide = prs.slides.add_slide(prs.slide_layouts[THANK_YOU_SLIDE_LAYOUT])
        thanks_slide.shapes.title.text = slides[-1].get('title', 'Thank You')

        if len(thanks_slide.placeholders) > 1:
            thanks_sub = thanks_slide.placeholders[1]
            thanks_sub.text = slides[-1].get('content', 'Thank you for your attention!')

        # Save final output
        output_filename = "generated_presentation.pptx"
        prs.save(output_filename)
        print("Presentation saved successfully")
        return output_filename

    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        traceback.print_exc()
        return ""

    
def extract_and_caption_pdf_elements(
    pdf_file_path: str,
    model_provider: str,
    model_name: str,
    api_key: str,
    output_dir: str = "./content/"
) -> list:
    """
    1) Extract figure images & possible table blocks using PyMuPDF.
    2) For tables: Summarize with a text LLM (HuggingFaceHub or ChatOpenAI).
    3) For images: Use a local pipeline("image-to-text") for BLIP or other captioning.
    """
    import os
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs("static", exist_ok=True)
    
    all_results = []
    
    try:
        import pymupdf as fitz
    except ImportError:
        raise ImportError("PyMuPDF is required for PDF processing. `pip install pymupdf`")
    
    try:
        doc = fitz.open(pdf_file_path)
        figure_counter = 1
        table_counter = 1
        
        ###############################################################################
        # 1) SET UP TEXT LLM (FOR SUMMARIZING TABLES)
        ###############################################################################
        if model_provider == "OpenAI":
            # Use ChatOpenAI for text-based tasks
            from langchain_openai import ChatOpenAI
            text_llm = ChatOpenAI(
                model=model_name,
                temperature=0.7,
                openai_api_key=api_key
            )
        else:
            # For text summarization with HF, use HuggingFaceHub (for Llama2, Flan, etc.)
            from langchain_community.llms import HuggingFaceHub
            text_llm = HuggingFaceHub(
                repo_id=model_name,
                task="text-generation",  
                huggingfacehub_api_token=api_key,
                model_kwargs={"temperature": 0.7}
            )
        
        ###############################################################################
        # 2) SET UP LOCAL IMAGE CAPTIONING PIPELINE 
        #    (DO NOT use HuggingFaceHub for BLIP/vit-gpt2 image-to-text)
        ###############################################################################
        # image_captioner = pipeline("image-to-text", model="Salesforce/blip-image-captioning-base")
        # image_captioner = get_image_captioner()
        max_retries = 3
        image_captioner = None
        
        for attempt in range(max_retries):
            try:
                image_captioner = get_image_captioner(model_name, api_key, model_provider)
                if image_captioner:
                    break
            except Exception as e:
                print(f"Attempt {attempt + 1} failed: {str(e)}")
                if attempt < max_retries - 1:
                    time.sleep(2)  # Wait before retry
                    
        if not image_captioner:
            print("Failed to initialize image captioner, will skip image captions")
            return []

        # If user selected "Salesforce/blip-image-captioning-large" or "nlpconnect/vit-gpt2-image-captioning"
        # we pass that model name to pipeline("image-to-text"). 
        # If your HF model is private/gated, you might also pass `use_auth_token=api_key`.
        # image_captioner = pipeline(
        #     task="image-to-text",
        #     model=model_name,    # e.g. "Salesforce/blip-image-captioning-large"
        #     revision="main",
        #     use_auth_token=api_key
        # )
        
        ###############################################################################
        # 3) PARSE THE PDF PAGES
        ###############################################################################
        for page_num in range(len(doc)):
            page = doc[page_num]
            page_rect = page.rect
            page_width = page_rect.width
            page_height = page_rect.height

            # ----------------------------------------------------------
            # PROCESS IMAGES ON THIS PAGE
            # ----------------------------------------------------------
            image_list = page.get_images(full=True)
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    if not base_image or not is_valid_figure(base_image):
                        continue

                    # Some quick heuristics to skip small images or page edges
                    valid_image = False
                    surrounding_text = ""
                    for img_rect in page.get_image_rects(xref):
                        if (img_rect.y0 < page_height * 0.1
                            or img_rect.y1 > page_height * 0.9
                            or ((img_rect.x1 - img_rect.x0) * (img_rect.y1 - img_rect.y0) < page_width * page_height * 0.05)
                        ):
                            continue
                        valid_image = True
                        # Optionally capture text near the figure
                        surrounding_text = get_surrounding_text(page, img_rect)
                        break

                    if not valid_image:
                        continue

                    # Save image to disk
                    ext = base_image["ext"].lower().replace("jpeg", "jpg")
                    image_filename = f"figure_{page_num}_{img_index}.{ext}"
                    image_path = os.path.join(output_dir, image_filename)
                    static_path = os.path.join("static", image_filename)

                    with open(image_path, "wb") as f, open(static_path, "wb") as sf:
                        f.write(base_image["image"])
                        sf.write(base_image["image"])

                    # CAPTION THIS IMAGE LOCALLY (not with HuggingFaceHub)
                    from PIL import Image
                    import io
                    raw_img = Image.open(io.BytesIO(base_image["image"])).convert("RGB")
                    if raw_img.size[0] < 10 or raw_img.size[1] < 10:
                        print(f"Skipping too small image on page {page_num}")
                        continue
                    # pipeline("image-to-text") returns a list of dict, e.g. [{"generated_text": "..."}]
                    try:
                        caption = caption_image(raw_img, image_captioner, model_provider)
                    except Exception as caption_err:
                        print(f"Caption generation failed: {str(caption_err)}")
                        caption = "Image caption unavailable"

                    all_results.append({
                        "type": "figure",
                        "file_path": image_path,
                        "static_path": static_path,
                        "caption": caption,
                        "figure_number": figure_counter
                    })
                    figure_counter += 1

                except Exception as e:
                    print(f"Error processing image on page {page_num}: {str(e)}")
                    continue

            # ----------------------------------------------------------
            # PROCESS TABLE-LIKE TEXT BLOCKS ON THIS PAGE
            # ----------------------------------------------------------
            text_blocks = page.get_text().split('\n\n')
            for block in text_blocks:
                if not is_likely_table(block):
                    continue
                
                try:
                    table_filename = f"table_{page_num}_{table_counter}.txt"
                    table_path = os.path.join(output_dir, table_filename)
                    static_table_path = os.path.join("static", table_filename)
                    
                    with open(table_path, "w", encoding="utf-8") as f1, open(static_table_path, "w", encoding="utf-8") as f2:
                        f1.write(block)
                        f2.write(block)
                    
                    # Summarize the table using the text LLM
                    from langchain_core.prompts import ChatPromptTemplate
                    from langchain_core.output_parsers import StrOutputParser

                    table_chain = (
                        ChatPromptTemplate.from_template("Summarize this table:\n{content}")
                        | text_llm
                        | StrOutputParser()
                    )
                    
                    caption = table_chain.invoke({"content": block})
                    
                    all_results.append({
                        "type": "table",
                        "file_path": table_path,
                        "static_path": static_table_path,
                        "caption": caption,
                        "figure_number": table_counter
                    })
                    table_counter += 1

                except Exception as e:
                    print(f"Error processing table: {str(e)}")
                    continue

        doc.close()
        return all_results
    
    except Exception as e:
        print(f"PDF processing error: {str(e)}")
        return []


    
def format_table_content(content: str) -> list:
    """Convert table text content into structured rows."""
    lines = [line.strip() for line in content.split('\n') if line.strip()]
    rows = []
    current_row = []
    
    for line in lines:
        # Split by common delimiters
        if '|' in line:
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
        elif '\t' in line:
            cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
        else:
            # For text that might be cell content
            cells = [line]
            
        if cells:
            current_row.extend(cells)
            
        # Start new row if we have enough cells or special markers
        if len(current_row) > 0 and (len(current_row) >= 3 or line.endswith('.')):
            rows.append(current_row)
            current_row = []
            
    # Add any remaining content
    if current_row:
        rows.append(current_row)
        
    return rows

def add_table_to_slide(slide, content: str, left: float, top: float, width: float) -> float:
    """Add a formatted table to the slide and return the bottom position."""
    try:
        # Format the content into rows
        rows = format_table_content(content)
        if not rows:
            return top
            
        # Determine number of columns
        max_cols = max(len(row) for row in rows)
        
        # Create table
        table_height = Inches(0.4) * len(rows)  # Estimate height
        table = slide.shapes.add_table(
            len(rows), max_cols,
            Inches(left),
            Inches(top),
            Inches(width),
            table_height
        ).table
        
        # Format table
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                if j < max_cols:  # Ensure we don't exceed table columns
                    table.cell(i, j).text = cell.strip()
                    # Format text
                    paragraph = table.cell(i, j).text_frame.paragraphs[0]
                    paragraph.font.size = Pt(10)
                    paragraph.font.name = "Calibri"
                    
                    # First row formatting
                    if i == 0:
                        paragraph.font.bold = True
                        table.cell(i, j).fill.solid()
                        table.cell(i, j).fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        # Auto-fit
        table.columns[0].width = Inches(width / max_cols)
        
        # Return the position below the table
        return top + (table_height / Inches(1)) + 0.2
        
    except Exception as e:
        print(f"Error adding table: {str(e)}")
        return top

async def add_formatted_table_element(slide, elem, current_y, text_llm):
    """
    Add a formatted table with improved spacing and styling.
    `current_y` is assumed to be in *inches*, so we do `Inches(current_y)` where needed.
    """
    try:
        table_width_in = 8.0
        left_in = 0.75

        # 1) read raw table text
        with open(elem['static_path'], 'r', encoding='utf-8') as f:
            raw_content = f.read()

        # 2) LLM to parse the table
        table_data = await process_table_text(raw_content, text_llm)
        if not table_data.headers and not table_data.rows:
            # No structured table found
            return current_y_in + 1.0

        # Estimate row heights
        num_rows = len(table_data.rows) + 1
        row_height_in = 0.3
        table_height_in = num_rows * row_height_in

        # Possibly move to new slide if not enough space
        slide = elem.get("slide_object", slide)  # or pass the slide reference from elsewhere
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN

        # Check vertical space
        # We'll do a simple function as an example:
        slide_width_in = slide.part.slide_width / 914_400.0
        slide_height_in = slide.part.slide_height / 914_400.0

        # if it doesn't fit, go to next slide
        slide, current_y_in = ensure_space_on_slide(slide.part, slide, current_y_in, table_height_in)

        # 3) Create the table shape
        table_shape = slide.shapes.add_table(
            rows=num_rows,
            cols=len(table_data.headers),
            left=Inches(left_in),
            top=Inches(current_y_in),
            width=Inches(table_width_in),
            height=Inches(table_height_in)
        )
        tbl = table_shape.table

        # 4) Fill headers
        for col_idx, header_text in enumerate(table_data.headers):
            cell = tbl.cell(0, col_idx)
            cell.text = header_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 230, 230)

        # 5) Fill rows
        for row_idx, row_obj in enumerate(table_data.rows, start=1):
            for col_idx, cell_text in enumerate(row_obj.cells):
                if col_idx < len(table_data.headers):
                    cell = tbl.cell(row_idx, col_idx)
                    cell.text = cell_text
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(10)

                    # Align numeric
                    if cell_text.replace('.', '', 1).replace('-', '', 1).isdigit():
                        paragraph.alignment = PP_ALIGN.RIGHT
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT

                    # Alternate row color
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

        # 6) Adjust column widths
        col_ratios = [1.5 if i == 0 else 1 for i in range(len(table_data.headers))]
        ratio_sum = sum(col_ratios)

        for i, col in enumerate(tbl.columns):
            col_width_emu = int(Inches(table_width_in) * (col_ratios[i] / ratio_sum))
            col.width = col_width_emu

        # 7) Add table caption below it
        caption_top_in = current_y_in + table_height_in + 0.2
        caption_box = slide.shapes.add_textbox(
            Inches(left_in),
            Inches(caption_top_in),
            Inches(table_width_in),
            Inches(0.4)
        )
        caption_para = caption_box.text_frame.add_paragraph()
        caption_para.text = f"Table {elem['figure_number']}: {elem['caption']}"
        caption_para.font.size = Pt(10)
        caption_para.font.italic = True
        caption_para.alignment = PP_ALIGN.CENTER

        return caption_top_in + 0.5

    except Exception as e:
        print(f"Error adding formatted table: {str(e)}")
        traceback.print_exc()
        return current_y + 3

# =================================
# Slide Placement Helpers
# =================================
def ensure_space_on_slide(prs, slide, current_y_in, shape_height_in):
    """
    If there's not enough vertical space left on the current slide, create a new slide.
    Return (slide, new_current_y_in).
    """
    slide_height_in = prs.slide_height / 914_400.0
    bottom_margin_in = 0.5  # for example

    # If shape won't fit:
    if (current_y_in + shape_height_in + bottom_margin_in) > slide_height_in:
        # Create new slide (using a standard content layout = 1 or whichever you prefer)
        new_slide = prs.slides.add_slide(prs.slide_layouts[1])
        # Reset to a top margin
        return (new_slide, 1.0)
    else:
        return (slide, current_y_in)

def add_image_within_bounds(prs, slide, file_path, current_y_in, left_in=1.0, desired_width_in=6.0, margin_bottom_in=0.2):
    """
    Add an image onto 'slide' at 'current_y_in' (in inches). 
    Auto-scale if it doesn't fit horizontally or vertically.
    Return new (slide, new_current_y_in).
    """
    from PIL import Image

    slide_width_in = prs.slide_width / 914_400.0
    slide_height_in = prs.slide_height / 914_400.0

    # Read actual image size in pixels
    img = Image.open(file_path)
    w_px, h_px = img.size
    aspect = w_px / float(h_px) if h_px != 0 else 1.0

    # Space left horizontally:
    right_margin_in = 1.0
    available_width_in = slide_width_in - (left_in + right_margin_in)

    # If desired_width_in is bigger than what's available, clamp it
    if desired_width_in > available_width_in:
        desired_width_in = available_width_in

    # Compute height from aspect ratio
    desired_height_in = desired_width_in / aspect

    # Space left vertically:
    bottom_margin_in = 0.5
    available_height_in = slide_height_in - (current_y_in + bottom_margin_in)

    # If the image is too tall, shrink it
    if desired_height_in > available_height_in:
        desired_height_in = available_height_in
        desired_width_in = desired_height_in * aspect

    # If there's still no space, push to new slide
    if desired_height_in < 0.5:
        # means there's basically no space, so let's create a new slide
        slide, current_y_in = ensure_space_on_slide(prs, slide, current_y_in, desired_height_in)
        # reset desired width
        desired_width_in = min(6.0, slide_width_in - 2.0)
        desired_height_in = desired_width_in / aspect

    # One final re-check for new slide space
    slide, current_y_in = ensure_space_on_slide(prs, slide, current_y_in, desired_height_in)

    # Add image
    left_emu = Inches(left_in)
    top_emu = Inches(current_y_in)
    width_emu = Inches(desired_width_in)
    height_emu = Inches(desired_height_in)

    pic = slide.shapes.add_picture(file_path, left_emu, top_emu, width=width_emu, height=height_emu)

    # Update 'current_y_in'
    current_y_in += desired_height_in + margin_bottom_in

    return slide, current_y_in


def get_llm(model_provider: str, model_name: str, api_key: str, temperature=0.7):
    """Factory function to get appropriate LLM"""
    max_retries = 3
    retry_delay = 2
    
    for attempt in range(max_retries):
        try:
            if model_provider == "OpenAI":
                return ChatOpenAI(
                    openai_api_key=api_key,
                    model_name=model_name,
                    temperature=temperature
                )
            else:
                from langchain_community.llms import HuggingFaceHub
                return HuggingFaceHub(
                    repo_id=model_name,
                    huggingfacehub_api_token=api_key,
                    model_kwargs={
                        "temperature": temperature,
                        "max_length": 512,
                        "num_return_sequences": 1
                    }
                )
        except Exception as e:
            print(f"Attempt {attempt + 1} failed: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay *= 2
            else:
                raise
    

def generate_image_caption(image_path: str, model_name: str, api_key: str = None) -> str:
    """Generate caption using either OpenAI or HF model"""
    from PIL import Image
    import base64
    
    def image_to_base64(image_path):
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode("utf-8")
    
    if "gpt-4" in model_name:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": [{
                    "type": "text",
                    "text": "Describe this image in 1 line for scientific presentation."
                }, {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{image_to_base64(image_path)}"}
                }]
            }],
            max_tokens=300
        )
        return response.choices[0].message.content
    else:
        # Use Hugging Face model
        processor = BlipProcessor.from_pretrained(model_name)
        model = BlipForConditionalGeneration.from_pretrained(model_name)
        
        raw_image = Image.open(image_path).convert('RGB')
        inputs = processor(raw_image, return_tensors="pt")
        out = model.generate(**inputs)
        return processor.decode(out[0], skip_special_tokens=True)
    



class TableRow(BaseModel):
    cells: List[str] = Field(default_factory=list)

class TableData(BaseModel):
    headers: List[str] = Field(default_factory=list)
    rows: List[TableRow] = Field(default_factory=list)

async def process_table_text(text: str, text_llm) -> TableData:
    """Process raw table text using LLM to structure it."""
    prompt = """
    Analyze and structure this technical data into a clear table format.
    Guidelines:
    1. Split multi-word headers into multiple columns
    2. Extract numeric values into separate columns
    3. Identify clear relationships between parameters and values
    4. Format as pipe-separated data with clear headers
    5. Keep units with their values
    
    For example, if input contains:
    "TOCNF was estimated to be 1.1 mmol g-1"
    
    Format as:
    Parameter | Value | Unit
    TOCNF | 1.1 | mmol g-1
    
    Text to analyze:
    {content}
    """
    
    try:
        table_chain = (
            ChatPromptTemplate.from_template(prompt)
            | text_llm
            | StrOutputParser()
        )
        structured_text = await table_chain.ainvoke({"content": text})
        
        table_data = TableData()
        lines = [line.strip() for line in structured_text.split('\n') if line.strip() and '|' in line]
        
        if lines:
            # Parse headers
            table_data.headers = [h.strip() for h in lines[0].split('|') if h.strip()]
            
            # Parse data rows
            for line in lines[1:]:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                if cells:
                    table_data.rows.append(TableRow(cells=cells))
                    
        return table_data
    except Exception as e:
        print(f"Error processing table text: {str(e)}")
        return TableData()

def validate_ecological_terms(text: str) -> str:
    """Enhance text with ecological term highlighting"""
    words = set(nltk.word_tokenize(text.lower()))
    missing_terms = ECOLOGICAL_TERMS - words
    
    if missing_terms:
        text += "\n\nPotential Ecological Considerations:\n- " + "\n- ".join(missing_terms)
    
    return text

def ecological_tfidf_analysis(text: str) -> dict:
    """Analyze text for ecological relevance"""
    vectorizer = TfidfVectorizer(stop_words=stopwords.words('english'))
    tfidf_matrix = vectorizer.fit_transform([text])
    feature_names = vectorizer.get_feature_names_out()
    
    eco_scores = {
        term: tfidf_matrix[0, vectorizer.vocabulary_[term]]
        for term in ECOLOGICAL_TERMS if term in vectorizer.vocabulary_
    }
    
    return {
        'top_terms': sorted(eco_scores.items(), key=lambda x: x[1], reverse=True)[:5],
        'missing_terms': [term for term in ECOLOGICAL_TERMS if term not in feature_names]
    }
