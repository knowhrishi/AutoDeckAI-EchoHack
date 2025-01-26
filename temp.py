import streamlit as st
import tempfile
import pdfplumber
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches
import json
import pymupdf as fitz
from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage

st.title("PDF to Smart Slide Deck Converter")
st.write("Upload a PDF to generate a PowerPoint presentation with images and tables!")

# Initialize session state
if 'slides' not in st.session_state:
    st.session_state.slides = None

# Get OpenAI API key
openai_api_key = st.text_input("Enter your OpenAI API key:", type="password")

def extract_pdf_content(pdf_path):
    """Extract text, images, and tables from PDF"""
    text = ""
    images = []
    tables = []
    
    # Extract text and images using PyMuPDF
    doc = fitz.open(pdf_path)
    for page_num, page in enumerate(doc):
        text += page.get_text()
        
        # Extract images
        img_list = page.get_images(full=True)
        for img_index, img in enumerate(img_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            img_ext = base_image["ext"]
            img_path = f"image_page{page_num+1}_num{img_index+1}.{img_ext}"
            with open(img_path, "wb") as img_file:
                img_file.write(image_bytes)
            images.append(img_path)
    
    # Extract tables using pdfplumber
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            table = page.extract_table()
            if table:
                df = pd.DataFrame(table[1:], columns=table[0])
                tables.append(df)
    
    return text, images, tables

def generate_slide_deck_structure(text, images_count, tables_count):
    """Use LLM to generate slide deck structure"""
    llm = ChatOpenAI(
        temperature=0.7,
        model="gpt-4",
        openai_api_key=openai_api_key
    )
    
    prompt = f"""
    Create a PowerPoint slide deck structure based on the following document content.
    The document contains {images_count} images and {tables_count} tables that should be included.
    
    Follow these rules:
    1. Create maximum 10 slides
    2. Use clear section headings
    3. Include key points as bullet lists
    4. Reference images/tables by index (Image 1, Table 2, etc.)
    5. Maintain logical flow
    
    Document Content:
    {text[:15000]}
    
    Respond in this JSON format:
    {{
        "slides": [
            {{
                "title": "Slide Title",
                "content": ["Bullet 1", "Bullet 2"],
                "image": 1,
                "table": 2
            }}
        ]
    }}
    """
    
    messages = [
        SystemMessage(content="You're a professional presentation designer. Create structured PowerPoint slides."),
        HumanMessage(content=prompt)
    ]
    
    response = llm(messages)
    return json.loads(response.content)

def create_powerpoint(slide_data, images, tables):
    """Generate PowerPoint from structured data"""
    prs = Presentation()
    
    for slide_info in slide_data['slides']:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_info.get('title', 'Slide Title')
        
        content_box = slide.placeholders[1]
        content_text = "\n".join(slide_info.get('content', []))
        
        # Handle images
        if 'image' in slide_info:
            img_idx = slide_info['image'] - 1
            if 0 <= img_idx < len(images):
                try:
                    slide.shapes.add_picture(images[img_idx], Inches(1), Inches(2), height=Inches(4))
                    content_text += f"\n[Image {slide_info['image']} inserted]"
                except Exception as e:
                    st.error(f"Error adding image: {e}")
        
        # Handle tables
        if 'table' in slide_info:
            table_idx = slide_info['table'] - 1
            if 0 <= table_idx < len(tables):
                try:
                    df = tables[table_idx]
                    rows, cols = df.shape
                    table_shape = slide.shapes.add_table(rows+1, cols, Inches(1), Inches(2), Inches(8), Inches(0.5*rows)).table
                    
                    # Add headers
                    for col_idx, col_name in enumerate(df.columns):
                        table_shape.cell(0, col_idx).text = str(col_name)
                    
                    # Add data
                    for row_idx in range(rows):
                        for col_idx in range(cols):
                            table_shape.cell(row_idx+1, col_idx).text = str(df.iloc[row_idx, col_idx])
                    
                    content_text += f"\n[Table {slide_info['table']} inserted]"
                except Exception as e:
                    st.error(f"Error adding table: {e}")
        
        content_box.text = content_text
    
    return prs

# File uploader and main logic remains the same
uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

if uploaded_file and openai_api_key:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        pdf_path = tmp_file.name
    
    with st.spinner("Extracting content from PDF..."):
        text, images, tables = extract_pdf_content(pdf_path)
    
    with st.spinner("Generating slide structure with AI..."):
        try:
            slide_structure = generate_slide_deck_structure(text, len(images), len(tables))
            st.session_state.slides = slide_structure
        except Exception as e:
            st.error(f"Error generating slides: {e}")
    
    if st.session_state.slides:
        st.subheader("Generated Slide Structure")
        st.json(st.session_state.slides)
        
        with st.spinner("Building PowerPoint..."):
            try:
                presentation = create_powerpoint(st.session_state.slides, images, tables)
                output_path = "generated_presentation.pptx"
                presentation.save(output_path)
                
                with open(output_path, "rb") as f:
                    st.download_button(
                        label="Download PowerPoint",
                        data=f,
                        file_name="generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            except Exception as e:
                st.error(f"Error creating PowerPoint: {e}")
    
    os.unlink(pdf_path)
    for img in images:
        if os.path.exists(img):
            os.remove(img)